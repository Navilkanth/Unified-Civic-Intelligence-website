import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def build_deck():
    prs = Presentation()
    
    # Use 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom color palette matching the web app styling
    navy = RGBColor(11, 31, 59)       # #0b1f3b
    teal = RGBColor(13, 107, 118)     # #0d6b76
    gold = RGBColor(201, 162, 39)     # #c9a227
    dark_gray = RGBColor(60, 60, 60)
    white = RGBColor(255, 255, 255)
    soft_bg = RGBColor(240, 244, 248)

    # 1. Slide Master & Custom Theme helper
    def apply_slide_bg(slide, rgb_color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color

    def add_title(slide, text, color=navy, size=40):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Georgia"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        return title_box

    # ──────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Dark Theme)
    # ──────────────────────────────────────────────────────────
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide1, navy)
    
    # Giant decorative shape text frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "UNIFIED CIVIC INTELLIGENCE"
    p.font.name = "Georgia"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = gold
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Civic Governance, Volunteer & Charity Platform"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.color.rgb = white
    p2.space_before = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Next-Generation React & Flask Ecosystem with Trust Ledger Integration"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(180, 200, 220)
    p3.space_before = Pt(40)

    # ──────────────────────────────────────────────────────────
    # SLIDE 2: Vision & Executive Summary (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide2, soft_bg)
    add_title(slide2, "Platform Vision & Value Proposition")
    
    left_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Core Pillars of Modern Governance"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = teal
    p.space_after = Pt(14)
    
    bullets = [
        "Empowers citizens with real-time feedback loops directly to administrators.",
        "Drives TVK Singapadai grassroots mobilization through a gamified point & badge engine.",
        "Establishes a zero-leakage, immutable Trust Ledger for public donations and welfare.",
        "Injects Heuristic NLP & AI signals for automatic complaint classification & analysis."
    ]
    for bullet in bullets:
        bp = tf_left.add_paragraph()
        bp.text = "• " + bullet
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Right Box - Key Stats / Visual Elements
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Why Unified Civic Intelligence?"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = navy
    p.space_after = Pt(14)
    
    reasons = [
        "100% Transparency: Every single donation, beneficiary allocation, and volunteer hour is logged to a clean public audit ledger.",
        "Interactive Engagement: Modern HSL glassmorphic design and responsive GIS map overlays excite and engage the digital-first generation.",
        "Scalable Architecture: Powered by a robust React frontend proxying JSON seamlessly to a multithreaded Flask core."
    ]
    for r in reasons:
        rp = tf_right.add_paragraph()
        rp.text = r
        rp.font.name = "Arial"
        rp.font.size = Pt(15)
        rp.font.color.rgb = dark_gray
        rp.space_before = Pt(12)

    # ──────────────────────────────────────────────────────────
    # SLIDE 3: Module Focus - Civic Governance (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide3, soft_bg)
    add_title(slide3, "Core Modules: Civic Governance Hub")
    
    col1 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col1 = col1.text_frame
    tf_col1.word_wrap = True
    
    p = tf_col1.paragraphs[0]
    p.text = "📋 Issues & AI Triage"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets1 = [
        "Automatic Category Detection: Infrastructure, negligence, corruption.",
        "Sentiment Analyzer: Flags urgent community pain points instantly.",
        "Anonymity Toggles: Encourages honest whistleblowing."
    ]
    for b in bullets1:
        bp = tf_col1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 2
    col2 = slide3.shapes.add_textbox(Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col2 = col2.text_frame
    tf_col2.word_wrap = True
    
    p = tf_col2.paragraphs[0]
    p.text = "💰 Fund & Works Tracking"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets2 = [
        "Contractor details, contact information, and sanctioned budgets.",
        "Real-time progress bars reflecting exact completion percentages.",
        "Interactive local Councillor office directories and hours."
    ]
    for b in bullets2:
        bp = tf_col2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 3
    col3 = slide3.shapes.add_textbox(Inches(8.85), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col3 = col3.text_frame
    tf_col3.word_wrap = True
    
    p = tf_col3.paragraphs[0]
    p.text = "🌍 Protection & Rescue"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets3 = [
        "Land Protection: Geo-coordinated tracking for encroachment reporting.",
        "Emergency Mode: Fast dispatch engine requiring zero authentication.",
        "Leaflet GIS Ward Overlays: Visualizing neighborhood analytics."
    ]
    for b in bullets3:
        bp = tf_col3.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # ──────────────────────────────────────────────────────────
    # SLIDE 4: Module Focus - TVK Singapadai Hub (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide4, soft_bg)
    add_title(slide4, "Core Modules: TVK Singapadai Volunteer Hub")
    
    left_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Grassroots Mobilization Engine"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = navy
    p.space_after = Pt(14)
    
    bullets = [
        "Live local RSS feeds highlighting ward-level news and announcements.",
        "District & constituency level Event managers for rallies, meetings, and donation drives.",
        "Actionable Task boards where coordinators assign direct duties (e.g. food prep, flood rescue).",
        "Point Reward system where completed tasks automatically recalculate volunteer standing."
    ]
    for b in bullets:
        bp = tf_left.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    right_box = slide4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Volunteer Gamification"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = gold
    p.space_after = Pt(14)
    
    reasons = [
        "Digital Membership: Successful registration generates a clean identity card equipped with custom ward QR code.",
        "Leaderboard System: Displays rank, name, ward, and points to promote healthy service-oriented competition.",
        "Digital Badges: Awards special achievements like 'Singapadai Starter' and 'Top Contributor' directly to profiles."
    ]
    for r in reasons:
        rp = tf_right.add_paragraph()
        rp.text = "✔ " + r
        rp.font.name = "Arial"
        rp.font.size = Pt(15)
        rp.font.color.rgb = dark_gray
        rp.space_before = Pt(12)

    # ──────────────────────────────────────────────────────────
    # SLIDE 5: Module Focus - Charity & Welfare Ledger (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide5, soft_bg)
    add_title(slide5, "Core Modules: Charity & Trust Ledger")
    
    col1 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col1 = col1.text_frame
    tf_col1.word_wrap = True
    
    p = tf_col1.paragraphs[0]
    p.text = "🤲 Welfare Requests"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets1 = [
        "Categorized support requests: Medical, food, and scholarship.",
        "AI Welfare Scheme Triage: Matches narratives to verified government programs.",
        "Structured documents upload and status reviews."
    ]
    for b in bullets1:
        bp = tf_col1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 2
    col2 = slide5.shapes.add_textbox(Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col2 = col2.text_frame
    tf_col2.word_wrap = True
    
    p = tf_col2.paragraphs[0]
    p.text = "📊 Transparency Hub"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets2 = [
        "Dynamic dashboard tracking total collections vs distributions.",
        "Sponsors Showcase displaying verified corporate contributions.",
        "Anonymous donation switches with direct ledger reference code generation."
    ]
    for b in bullets2:
        bp = tf_col2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 3
    col3 = slide5.shapes.add_textbox(Inches(8.85), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col3 = col3.text_frame
    tf_col3.word_wrap = True
    
    p = tf_col3.paragraphs[0]
    p.text = "📒 Trust Ledger"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = teal
    
    bullets3 = [
        "Immutable ledger entry tracking: donation, distribution, service.",
        "Automatic volunteer service logging to ledger upon task signoff.",
        "Instant filtering by entry type to guarantee public audit integrity."
    ]
    for b in bullets3:
        bp = tf_col3.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # ──────────────────────────────────────────────────────────
    # SLIDE 6: AI Engine & Smart Capabilities (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide6, soft_bg)
    add_title(slide6, "AI Engine & Smart Capabilities")
    
    left_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "The UCI AI Brain"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = navy
    p.space_after = Pt(14)
    
    bullets = [
        "Heuristic Natural Language Processing engine analyzes user complaint descriptions.",
        "Predicts the best category fit and extracts the underlying public sentiment score (-1 to +1).",
        "Translates long citizen welfare descriptions into targeted, matching government schemes.",
        "Highlights elevated stress levels in specific wards to trigger proactive councillor reviews."
    ]
    for b in bullets:
        bp = tf_left.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    right_box = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Interactive Conversational Agent"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = teal
    p.space_after = Pt(14)
    
    reasons = [
        "Rule-Based Expert System: Fully maps to 10 distinct civic query categories.",
        "Quick Prompts: Dynamic clicks guide citizens to view project progress, submit land issues, or search for welfare schemes.",
        "Sentiment-Aware fallback responses support and guide frustrated users constructively."
    ]
    for r in reasons:
        rp = tf_right.add_paragraph()
        rp.text = "✔ " + r
        rp.font.name = "Arial"
        rp.font.size = Pt(15)
        rp.font.color.rgb = dark_gray
        rp.space_before = Pt(12)

    # ──────────────────────────────────────────────────────────
    # SLIDE 7: System Architecture (Light Theme)
    # ──────────────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide7, soft_bg)
    add_title(slide7, "System Architecture & Security")
    
    col1 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col1 = col1.text_frame
    tf_col1.word_wrap = True
    
    p = tf_col1.paragraphs[0]
    p.text = "💻 React SPA (Frontend)"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = navy
    
    bullets1 = [
        "Built with React 18, Vite 8, and custom Tailwind CSS v3 styling.",
        "AuthContext with automated JWT tokens in local storage.",
        "State management with transparent CORS api wrappers."
    ]
    for b in bullets1:
        bp = tf_col1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 2
    col2 = slide7.shapes.add_textbox(Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col2 = col2.text_frame
    tf_col2.word_wrap = True
    
    p = tf_col2.paragraphs[0]
    p.text = "⚙ Flask Core (Backend)"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = navy
    
    bullets2 = [
        "REST API blueprint structure routing requests to specific modules.",
        "Socket.io server mapping real-time notifications to active room lists.",
        "Multithreaded SQLite / PostgreSQL engine using SQLAlchemy."
    ]
    for b in bullets2:
        bp = tf_col2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # Column 3
    col3 = slide7.shapes.add_textbox(Inches(8.85), Inches(1.8), Inches(3.6), Inches(4.8))
    tf_col3 = col3.text_frame
    tf_col3.word_wrap = True
    
    p = tf_col3.paragraphs[0]
    p.text = "🔒 Security & Compliance"
    p.font.name = "Georgia"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = navy
    
    bullets3 = [
        "Password Hashing: Advanced PBKDF2/SHA256 hash checks.",
        "Role-Based Access Control: Decorators protect admin and volunteer endpoints.",
        "Audit Logging: Track every complaint, donation, and login event."
    ]
    for b in bullets3:
        bp = tf_col3.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = dark_gray
        bp.space_before = Pt(10)

    # ──────────────────────────────────────────────────────────
    # SLIDE 8: Project Deliverables & Future (Dark Theme)
    # ──────────────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide8, navy)
    add_title(slide8, "Project Status & Deliverables", gold, 40)
    
    left_box = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Completed Milestones"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white
    p.space_after = Pt(14)
    
    bullets = [
        "100% Migration: Fully replaced Jinga2 UI with active React components.",
        "Stabilized Core: Reloader issues solved, switched safely to multithreaded mode.",
        "Rich UI: Visual excellence backed by curated HSL CSS theme elements.",
        "API Stability: 100% CORS-friendly JSON APIs verified and running."
    ]
    for b in bullets:
        bp = tf_left.add_paragraph()
        bp.text = "✔ " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(15)
        bp.font.color.rgb = RGBColor(220, 230, 240)
        bp.space_before = Pt(10)

    right_box = slide8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Next Scaling Steps"
    p.font.name = "Georgia"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = gold
    p.space_after = Pt(14)
    
    steps = [
        "Geographic expansion of Leaflet boundary overlays to all 200 Chennai Wards.",
        "Integrate live payment gateway wrappers (Razorpay/Stripe) into the Charity Hub.",
        "Mobile App adaptation leveraging modern React Native wrappers."
    ]
    for s in steps:
        sp = tf_right.add_paragraph()
        sp.text = "• " + s
        sp.font.name = "Arial"
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(220, 230, 240)
        sp.space_before = Pt(12)

    prs.save("Unified_Civic_Ecosystem_Platform.pptx")
    print("Presentation generated successfully!")

if __name__ == "__main__":
    build_deck()
